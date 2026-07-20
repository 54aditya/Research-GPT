from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.core.config import get_settings

settings = get_settings()


class PresentationAgent:
    def create_presentation(self, topic: str, slides: list[dict[str, object]], presentation_style: str) -> str:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = topic
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        
        subtitle_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = f"{presentation_style.title()} Research Presentation"
        subtitle_p.font.size = Pt(28)

        # Agenda Slide
        agenda_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        agenda_title = agenda_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        agenda_title_frame = agenda_title.text_frame
        agenda_title_p = agenda_title_frame.paragraphs[0]
        agenda_title_p.text = "Agenda"
        agenda_title_p.font.size = Pt(44)
        agenda_title_p.font.bold = True
        
        agenda_content = agenda_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5.5))
        agenda_frame = agenda_content.text_frame
        agenda_frame.word_wrap = True
        for idx, slide_data in enumerate(slides[:12], 1):  # Show all agenda items up to 12
            if idx > 1:
                agenda_frame.add_paragraph()
            p = agenda_frame.paragraphs[idx - 1]
            p.text = f"{idx}. {slide_data.get('title', f'Point {idx}')}"
            p.font.size = Pt(18)
            p.level = 0

        # Content Slides
        for slide_idx, slide_data in enumerate(slides, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Slide number and title
            slide_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            slide_title_frame = slide_title_box.text_frame
            slide_title_frame.word_wrap = True
            slide_title_p = slide_title_frame.paragraphs[0]
            slide_title_p.text = f"{slide_data.get('title', f'Slide {slide_idx}')}"
            slide_title_p.font.size = Pt(40)
            slide_title_p.font.bold = True
            
            # Content
            content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            content_text = slide_data.get("content", "")
            if content_text:
                # Split content into individual paragraphs/lines to parse list items
                lines = [line.strip() for line in str(content_text).split('\n') if line.strip()]
                
                # Determine font size dynamically based on length to prevent text overflow
                total_len = len(content_text)
                if total_len > 1200:
                    font_size = Pt(10)
                elif total_len > 1000:
                    font_size = Pt(11)
                elif total_len > 800:
                    font_size = Pt(12)
                elif total_len > 600:
                    font_size = Pt(13)
                elif total_len > 400:
                    font_size = Pt(14)
                else:
                    font_size = Pt(16)
                
                first = True
                for line in lines:
                    level = 0
                    # Identify lists starting with standard bullet indicators
                    if line.startswith(('-', '*', '•')):
                        line = line.lstrip('-*•').strip()
                        level = 1
                    elif line.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')) and len(line) > 2 and line[2] == ' ':
                        line = line[2:].strip()
                        level = 1

                    if first:
                        p = content_frame.paragraphs[0]
                        first = False
                    else:
                        p = content_frame.add_paragraph()

                    p.text = line
                    p.font.size = font_size
                    p.space_before = Pt(4)
                    p.space_after = Pt(4)
                    p.level = level
            else:
                # Default content if empty
                p = content_frame.paragraphs[0]
                p.text = "Research content and key findings"
                p.font.size = Pt(16)
            
            # Slide number footer
            footer_box = slide.shapes.add_textbox(Inches(9), Inches(7.2), Inches(0.8), Inches(0.3))
            footer_frame = footer_box.text_frame
            footer_p = footer_frame.paragraphs[0]
            footer_p.text = f"{slide_idx}"
            footer_p.font.size = Pt(12)
            footer_p.alignment = PP_ALIGN.RIGHT

        output_dir = Path(settings.upload_dir) / "ppts"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{topic.replace(' ', '_')}.pptx"
        prs.save(output_path)
        return str(output_path.resolve())
